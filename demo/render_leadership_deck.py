#!/usr/bin/env python3
"""Render the S7 leadership overview deck (docs/s7-leadership-overview.pptx).

Authoring tooling, not runtime — same discipline as demo/render_pdf.py: this
script and its python-pptx dependency are not required to run the demo, only
to regenerate this one deliverable, so it stays outside hard rule 4's
portability requirement.

Content is the same eight slides as the original build (title, one-sentence
pitch, pipeline flow, five gates, traceability, ten building blocks, next
steps, takeaway) — this version only changes the *skin*, to match the S3
console's ported design system (theme.css / apps/control/web/src/theme.css):
warm near-black ink, cream surface, crimson accent, square corners
everywhere, Source Sans typography. See CLAUDE.md hard rule 2 — MapleSure
only, the end client is never named.

    .venv/bin/python3 demo/render_leadership_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "s7-leadership-overview.pptx"

# --- palette, ported from apps/control/web/src/theme.css ------------------
INK = RGBColor(0x36, 0x36, 0x2F)
INK_SOFT = RGBColor(0x52, 0x51, 0x46)
BG = RGBColor(0xF0, 0xF0, 0xEC)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD2, 0xD2, 0xCF)
LINE_STRONG = RGBColor(0x85, 0x84, 0x7A)
ACCENT = RGBColor(0xA2, 0x0A, 0x29)
ACCENT_DARK = RGBColor(0x7A, 0x08, 0x20)
ACCENT_SOFT = RGBColor(0xF7, 0xE9, 0xEC)
GOLD = RGBColor(0xB4, 0x88, 0x25)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Source Sans Pro"

prs = Presentation()
prs.slide_width = Emu(12191695)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, bg, line=None)
    return s


def rect(s, l, t, w, h, fill, line=LINE, line_w=Pt(1)):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def arrow(s, l, t, w, h, fill=LINE_STRONG):
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs: list of (text, size_pt, bold, color) tuples, one per paragraph."""
    box = s.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (txt, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return box


def kicker_rule(s, l, t):
    rect(s, l, t, Emu(732536), Emu(38100), ACCENT, line=None)


def title_block(s, heading, sub=None):
    text(s, 548640, 320040, 11064240, 685800, [(heading, 34, True, INK)])
    kicker_rule(s, 594360, 987552)
    if sub:
        text(s, 548640, 1115568, 11064240, 457200, [(sub, 16, False, INK_SOFT)])


# --- Slide 1 — title -------------------------------------------------------
s = slide(bg=INK)
text(s, 914400, 2331720, 10332720, 685800, [("AI-Assisted Delivery", 54, True, WHITE)])
text(s, 914400, 3017520, 10332720, 685800,
     [("From business requirement to production release — with humans in control", 24, False, RGBColor(0xE3, 0xC9, 0xCE))])
rect(s, 914400, 5750000, 2011680, Emu(38100), ACCENT, line=None)
text(s, 914400, 5897880, 10332720, 457200, [("S7 · Delivery Scope · August 2026", 15, False, RGBColor(0xB8, 0xA9, 0xA0))])

# --- Slide 2 — the idea in one sentence -------------------------------------
s = slide()
title_block(s, "The idea, in one sentence")
text(s, 822960, 2468880, 10515600, 2286000, [
    ("AI does the heavy lifting of delivery —", 32, True, INK),
    ("analysis, design, stories, code, tests, documentation —", 26, False, ACCENT_DARK),
    ("and nothing reaches production without passing our checks.", 32, True, INK),
])
text(s, 1463040, 5029200, 9235440, 731520,
     [("Not a black box. Every step is gated, logged, and traceable.", 19, False, INK_SOFT)])

# --- Slide 3 — pipeline flow -------------------------------------------------
s = slide()
title_block(s, "How work will flow", "One pipeline, start to finish")
STEPS = [
    ("Business\nrequirement", ACCENT),
    ("AI\nassessment", ACCENT),
    ("AI design\n(diagrams)", ACCENT),
    ("HUMAN\nREVIEW", GOLD),
    ("User\nstories", ACCENT),
    ("Build &\ntest", ACCENT),
    ("HUMAN\nAPPROVAL", GOLD),
    ("Release", RGBColor(0x12, 0x5C, 0x43)),
]
x = 502920
box_w, box_h, gap, arrow_w = 1298448, 960120, 237744, 237744
for i, (label, color) in enumerate(STEPS):
    rect(s, x, 2423160, box_w, box_h, color, line=None)
    text(s, x, 2423160, box_w, box_h, [(label, 13, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(STEPS) - 1:
        arrow(s, x + box_w, 2779776, arrow_w, 256032)
    x += box_w + gap
text(s, 822960, 3977639, 10515600, 2011680, [
    ("•  AI prepares every deliverable: assessment, design diagrams, stories, code, tests, docs.", 19, False, INK_SOFT),
    ("•  The gold boxes are decision points — a person reviews and approves before work continues.", 19, False, INK_SOFT),
    ("•  If the reviewer rejects, the pipeline stops. That is by design, not a limitation.", 19, False, INK_SOFT),
])

# --- Slide 4 — five gates ----------------------------------------------------
s = slide()
title_block(s, "Nothing advances without passing a gate", "Five checkpoints from intake to release")
GATES = [
    ("Gates 0–2", "Completeness checks",
     "Is every requirement mapped to a story? Is every story testable? Automated checks — before any human spends time reviewing."),
    ("Gate 3", "Independent review",
     "A second, separate AI reviews the first one's work against the design. Critical gaps block progress. No phase approves its own work."),
    ("Gate 4", "Human approval",
     "A person makes the final call before release. The system cannot ship on its own."),
]
y = 1920240
for label, title_txt, body in GATES:
    rect(s, 731520, y, 2103120, 1005840, ACCENT, line=None)
    text(s, 731520, y, 2103120, 1005840, [(label, 18, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3200400, y - 18288, 8321040, 1143000, [
        (title_txt, 20, True, ACCENT_DARK),
        (body, 16, False, INK_SOFT),
    ])
    y += 1371600
rect(s, 731520, 5989320, 10789920, 502920, ACCENT_SOFT, line=ACCENT)
text(s, 731520, 5989320, 10789920, 502920,
     [("The gates are real: a rejection stops the pipeline. We will demonstrate that live.", 16, True, ACCENT_DARK)],
     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# --- Slide 5 — traceability ---------------------------------------------------
s = slide()
title_block(s, "Every output can be traced back", "If something is wrong, we can find out why in minutes")
CARDS = [
    ("Provenance ledger", "A tamper-evident record of every artifact: who made it, when, and from what input."),
    ("Traceability matrix", "Requirement → design → story → code → test → release, linked end to end. A defect traces backward in one lookup."),
    ("Activity log", "Every AI session logged: what it did, how long it took, what it produced. Shows real velocity and real bottlenecks."),
    ("Staleness detection", "If a design changes, everything built on it is flagged out-of-date and must be refreshed before release."),
]
positions = [(731520, 1874519), (6309360, 1874519), (731520, 3977639), (6309360, 3977639)]
for (title_txt, body), (l, t) in zip(CARDS, positions):
    rect(s, l, t, 5166360, 1874519, SURFACE, line=LINE)
    rect(s, l, t, Emu(38100), 1874519, ACCENT, line=None)
    text(s, l + 274320, t + 228600, 4700000, 1420000, [
        (title_txt, 20, True, INK),
        (body, 15, False, INK_SOFT),
    ])

# --- Slide 6 — ten building blocks --------------------------------------------
s = slide()
title_block(s, "The ten building blocks", "Grouped by what they do for us")
GROUPS = [
    ("CONTROL", "Keeps humans in charge",
     ["Gated pipeline (5 gates)", "Completeness checks", "Independent review", "Staleness detection"], ACCENT),
    ("TRUST & AUDIT", "Proves what happened",
     ["Provenance ledger", "Traceability matrix", "Activity log"], GOLD),
    ("FOUNDATION", "Makes it repeatable",
     ["Four-layer architecture", "Story quality standards", "Change management"], RGBColor(0x12, 0x5C, 0x43)),
]
x = 731520
for label, sub, items, color in GROUPS:
    rect(s, x, 1874519, 3429000, 3977639, SURFACE, line=LINE)
    rect(s, x + 228600, 2148840, 2971800, 548640, color, line=None)
    text(s, x + 228600, 2148840, 2971800, 548640, [(label, 15, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 320040, 2880360, 2880360, 457200, [(sub, 14, False, INK_SOFT)])
    text(s, x + 320040, 3429000, 2926080, 2286000,
         [(f"•  {item}", 15, False, INK) for item in items])
    x += 3657600
text(s, 731520, 6126480, 10698480, 548640,
     [("Ten features. Three jobs. One governed pipeline.", 17, True, INK)], align=PP_ALIGN.CENTER)

# --- Slide 7 — what happens next -----------------------------------------------
s = slide()
title_block(s, "What happens next")
text(s, 822960, 1737360, 10515600, 3291840, [
    ("•  We build these ten capabilities point by point, in priority order.", 21, False, INK),
    ("•  Every stage ends with something that runs — a working demo, not a slide.", 21, False, INK),
    ("•  Anything simulated is clearly labelled as simulated. What you see live is real.", 21, False, INK),
    ("•  Human checkpoints stay on until a workflow has proven itself — autonomy is earned, not assumed.", 21, False, INK),
])
rect(s, 822960, 5257800, 10515600, 868680, ACCENT, line=None)
text(s, 822960, 5257800, 10515600, 868680,
     [("You will see the pipeline run end to end, watch a gate block bad work, and trace an output back to its requirement.", 17, True, WHITE)],
     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# --- Slide 8 — takeaway ---------------------------------------------------------
s = slide(bg=INK)
text(s, 1097280, 2103120, 9966960, 457200, [("THE TAKEAWAY", 18, True, RGBColor(0xE3, 0xC9, 0xCE))])
text(s, 1097280, 2560320, 9966960, 914400, [("Speed from AI.", 44, True, WHITE)])
text(s, 1097280, 3200400, 9966960, 914400, [("Confidence from governance.", 44, True, WHITE)])
rect(s, 1097280, 4200000, 2011680, Emu(38100), ACCENT, line=None)
text(s, 1097280, 4400000, 9966960, 731520,
     [("Every artifact gated, logged, and traceable — and a human always holds the final approval.", 20, False, RGBColor(0xE3, 0xC9, 0xCE))])

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"{OUT}  ({OUT.stat().st_size / 1024:.0f} KB)")
